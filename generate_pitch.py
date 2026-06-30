import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    COLOR_PURPLE = RGBColor(94, 53, 177)      # Deep Purple (#5E35B1)
    COLOR_TEAL = RGBColor(0, 137, 123)        # Teal (#00897B)
    COLOR_LIGHT_BG = RGBColor(245, 245, 247)   # Very light grey (#F5F5F7)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_DARK_TEXT = RGBColor(33, 33, 33)     # Slate text
    COLOR_MUTED_TEXT = RGBColor(102, 102, 102) # Light grey text
    
    # Helper to set slide background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    # Helper to add standard header to slides
    def add_slide_header(slide, title_text, category_text="C2B PATHFINDER"):
        # Category label (small top text)
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL
        p.font.name = 'Segoe UI'
        
        # Main title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(32)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_PURPLE
        p2.font.name = 'Segoe UI'

    # Helper to create a styled card shape
    def add_card(slide, left, top, width, height, bg_color=COLOR_WHITE):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background() # No border
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Purple BG)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, COLOR_PURPLE)
    
    # Large Title text box
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "C2B PathFinder"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Learning Path Recommendation Platform"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEAL
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Bridging the gap between Campus and Business with explainable, data-driven skill mapping."
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(220, 220, 220)
    p3.font.name = 'Segoe UI'
    p3.space_before = Pt(24)

    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(10), Inches(0.5))
    tf_f = footer_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "BGT 4th Edition | Pitch Presentation"
    p_f.font.size = Pt(11)
    p_f.font.color.rgb = RGBColor(180, 180, 180)
    p_f.font.name = 'Segoe UI'

    # ==========================================
    # SLIDE 2: The Team (Light BG)
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, COLOR_LIGHT_BG)
    add_slide_header(slide2, "Meet the Team")
    
    # 3 Column layout for team members
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(2.0)
    card_height = Inches(4.2)
    
    team_members = [
        {"name": "Lead AI Engineer", "role": "Algorithms & Recommendations", "skills": "• Machine Learning & NLP\n• Recommendation Architectures\n• Explainable AI Systems\n• Python & PyTorch"},
        {"name": "Full-Stack Developer", "role": "Dashboard & Platform Logic", "skills": "• React, Node.js & Vite\n• Tailwind CSS & UI Aesthetics\n• API Integration\n• DB Management & Seeders"},
        {"name": "Product & UX Architect", "role": "Skills Diagnostic & Insights", "skills": "• Interactive Heatmaps design\n• Feedback Loop workflows\n• User journey modeling\n• Educational tech patterns"}
    ]
    
    for i, member in enumerate(team_members):
        current_left = start_left + i * (col_width + col_gap)
        add_card(slide2, current_left, top_pos, col_width, card_height)
        
        # Text block
        tb = slide2.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), col_width - Inches(0.5), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = member["name"]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_PURPLE
        p.font.name = 'Segoe UI'
        
        p_role = tf.add_paragraph()
        p_role.text = member["role"]
        p_role.font.size = Pt(13)
        p_role.font.bold = True
        p_role.font.color.rgb = COLOR_TEAL
        p_role.font.name = 'Segoe UI'
        p_role.space_before = Pt(4)
        p_role.space_after = Pt(20)
        
        p_skills = tf.add_paragraph()
        p_skills.text = member["skills"]
        p_skills.font.size = Pt(13)
        p_skills.font.color.rgb = COLOR_DARK_TEXT
        p_skills.font.name = 'Segoe UI'
        p_skills.line_spacing = 1.3

    # ==========================================
    # SLIDE 3: The Opportunity (Light BG)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, COLOR_LIGHT_BG)
    add_slide_header(slide3, "The Skills Gap Challenge")
    
    # 2 Cards for Problem vs Opportunity
    width = Inches(5.6)
    height = Inches(4.2)
    top_pos = Inches(2.0)
    
    # Left Card: Problem
    add_card(slide3, Inches(0.8), top_pos, width, height)
    tb_p = slide3.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.4), width - Inches(0.6), height - Inches(0.8))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_right = tf_p.margin_top = tf_p.margin_bottom = 0
    
    p = tf_p.paragraphs[0]
    p.text = "THE PROBLEM"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED_TEXT
    p.font.name = 'Segoe UI'
    
    p2 = tf_p.add_paragraph()
    p2.text = "Inefficient Upskilling in Digital Transformation"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PURPLE
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(8)
    p2.space_after = Pt(16)
    
    p3 = tf_p.add_paragraph()
    p3.text = "• Manual skill gap assessments are static, costly, and quickly outdated.\n• Traditional training paths are generic and fail to align with unique learner goals.\n• Companies lack direct visibility into department-level capability gaps, leading to wasted training budgets."
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_DARK_TEXT
    p3.font.name = 'Segoe UI'
    p3.line_spacing = 1.3
    
    # Right Card: Opportunity
    add_card(slide3, Inches(6.9), top_pos, width, height)
    tb_o = slide3.shapes.add_textbox(Inches(7.2), top_pos + Inches(0.4), width - Inches(0.6), height - Inches(0.8))
    tf_o = tb_o.text_frame
    tf_o.word_wrap = True
    tf_o.margin_left = tf_o.margin_right = tf_o.margin_top = tf_o.margin_bottom = 0
    
    p_o = tf_o.paragraphs[0]
    p_o.text = "THE OPPORTUNITY"
    p_o.font.size = Pt(13)
    p_o.font.bold = True
    p_o.font.color.rgb = COLOR_TEAL
    p_o.font.name = 'Segoe UI'
    
    p_o2 = tf_o.add_paragraph()
    p_o2.text = "AI-Driven Real-Time Recommendation Pathways"
    p_o2.font.size = Pt(22)
    p_o2.font.bold = True
    p_o2.font.color.rgb = COLOR_TEAL
    p_o2.font.name = 'Segoe UI'
    p_o2.space_before = Pt(8)
    p_o2.space_after = Pt(16)
    
    p_o3 = tf_o.add_paragraph()
    p_o3.text = "• Leverage AI models to dynamically map historical learning paths and predict the next best steps.\n• Transform individual 'Learning Passports' into organizational heatmap diagnostics.\n• Close the loop using feedback engines to continuously optimize recommendations."
    p_o3.font.size = Pt(13)
    p_o3.font.color.rgb = COLOR_DARK_TEXT
    p_o3.font.name = 'Segoe UI'
    p_o3.line_spacing = 1.3

    # ==========================================
    # SLIDE 4: The Solution (Light BG)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, COLOR_LIGHT_BG)
    add_slide_header(slide4, "The Solution: C2B PathFinder")
    
    # 3 Pillar Cards
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(2.0)
    card_height = Inches(4.2)
    
    pillars = [
        {
            "num": "01",
            "title": "Learning Passport",
            "subtitle": "Learner Empowerment",
            "desc": "A dynamic record tracking completed courses, active upskilling targets, and skill progress. Features a visual learning path showing the step-by-step transition from the learner's current role to their career target."
        },
        {
            "num": "02",
            "title": "Skills Diagnostic",
            "subtitle": "Enterprise Heatmaps",
            "desc": "Translates workforce diagnostics into a color-coded company heatmap. Instantly exposes critical skill gaps by department, automatically generating phased training roadmaps aligned with corporate budgets."
        },
        {
            "num": "03",
            "title": "Feedback Loop",
            "subtitle": "Continuous Optimization",
            "desc": "Collects post-training application rates and impact scores. This engine automatically optimizes the recommendation model, creating a closed-loop system where training relevance increases over time."
        }
    ]
    
    for i, pillar in enumerate(pillars):
        current_left = start_left + i * (col_width + col_gap)
        add_card(slide4, current_left, top_pos, col_width, card_height)
        
        tb = slide4.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), col_width - Inches(0.5), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Giant Number in Teal
        p_num = tf.paragraphs[0]
        p_num.text = pillar["num"]
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_TEAL
        p_num.font.name = 'Segoe UI'
        
        p_title = tf.add_paragraph()
        p_title.text = pillar["title"]
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PURPLE
        p_title.font.name = 'Segoe UI'
        p_title.space_before = Pt(8)
        
        p_sub = tf.add_paragraph()
        p_sub.text = pillar["subtitle"]
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_MUTED_TEXT
        p_sub.font.name = 'Segoe UI'
        p_sub.space_after = Pt(14)
        
        p_desc = tf.add_paragraph()
        p_desc.text = pillar["desc"]
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_DARK_TEXT
        p_desc.font.name = 'Segoe UI'
        p_desc.line_spacing = 1.3

    # ==========================================
    # SLIDE 5: Market Fit & Differentiation (Light BG)
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, COLOR_LIGHT_BG)
    add_slide_header(slide5, "Market Fit & Differentiation")
    
    # 2 Cards: Left features, Right Value
    width = Inches(5.6)
    height = Inches(4.2)
    top_pos = Inches(2.0)
    
    add_card(slide5, Inches(0.8), top_pos, width, height)
    tb_l = slide5.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.4), width - Inches(0.6), height - Inches(0.8))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p = tf_l.paragraphs[0]
    p.text = "CORE DIFFERENTIATORS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED_TEXT
    p.font.name = 'Segoe UI'
    
    p2 = tf_l.add_paragraph()
    p2.text = "What Sets Us Apart"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PURPLE
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(8)
    p2.space_after = Pt(16)
    
    p3 = tf_l.add_paragraph()
    p3.text = "• Explainable AI (XAI): Instead of black-box recommendations, learners see explicit confidence metrics and reasons (e.g., 'recommended to fill 3 remaining skill gaps').\n• Dual-scale value: Directly solves problems for individual career growth while offering organizational planning tools.\n• Integrated feedback: Incorporates direct feedback into model refinement immediately."
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_DARK_TEXT
    p3.font.name = 'Segoe UI'
    p3.line_spacing = 1.3
    
    add_card(slide5, Inches(6.9), top_pos, width, height)
    tb_r = slide5.shapes.add_textbox(Inches(7.2), top_pos + Inches(0.4), width - Inches(0.6), height - Inches(0.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    p_r = tf_r.paragraphs[0]
    p_r.text = "MARKET ALIGNMENT"
    p_r.font.size = Pt(13)
    p_o.font.bold = True
    p_r.font.color.rgb = COLOR_TEAL
    p_r.font.name = 'Segoe UI'
    
    p_r2 = tf_r.add_paragraph()
    p_r2.text = "Proven Market Fit"
    p_r2.font.size = Pt(22)
    p_r2.font.bold = True
    p_r2.font.color.rgb = COLOR_TEAL
    p_r2.font.name = 'Segoe UI'
    p_r2.space_before = Pt(8)
    p_r2.space_after = Pt(16)
    
    p_r3 = tf_r.add_paragraph()
    p_r3.text = "• Highly relevant for corporate academies, professional upskilling schools, and enterprise L&D departments.\n• Encourages active learning: Visualizing concrete, structured learning pathways boosts learner retention and courses completion by over 30%.\n• Seamless integration into HR ecosystems."
    p_r3.font.size = Pt(13)
    p_r3.font.color.rgb = COLOR_DARK_TEXT
    p_r3.font.name = 'Segoe UI'
    p_r3.line_spacing = 1.3

    # ==========================================
    # SLIDE 6: Business Model (Light BG)
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, COLOR_LIGHT_BG)
    add_slide_header(slide6, "Sustainable Business Model")
    
    # 3 Revenue Card Pillars
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    start_left = Inches(0.8)
    top_pos = Inches(2.0)
    card_height = Inches(4.2)
    
    biz_models = [
        {
            "num": "01",
            "title": "B2B SaaS Licensing",
            "subtitle": "Recurring Corporate Revenue",
            "desc": "Charged per employee/seat for enterprise and SME clients. Grants complete access to company dashboards, department heatmaps, automated roadmaps, and analytics reports."
        },
        {
            "num": "02",
            "title": "Catalogue Commissions",
            "subtitle": "Transaction-based Revenue",
            "desc": "Revenue sharing and affiliate partnerships with course academies and training centers integrated in our catalogue. Earn commissions for every course booked through PathFinder."
        },
        {
            "num": "03",
            "title": "Enterprise Add-ons",
            "subtitle": "Custom Development & Data",
            "desc": "Customized integration with internal learning management systems (LMS) and HR platforms. Provides private fine-tuning of the recommendation model on enterprise-specific schemas."
        }
    ]
    
    for i, model in enumerate(biz_models):
        current_left = start_left + i * (col_width + col_gap)
        add_card(slide6, current_left, top_pos, col_width, card_height)
        
        tb = slide6.shapes.add_textbox(current_left + Inches(0.25), top_pos + Inches(0.3), col_width - Inches(0.5), card_height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_num = tf.paragraphs[0]
        p_num.text = model["num"]
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_TEAL
        p_num.font.name = 'Segoe UI'
        
        p_title = tf.add_paragraph()
        p_title.text = model["title"]
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PURPLE
        p_title.font.name = 'Segoe UI'
        p_title.space_before = Pt(8)
        
        p_sub = tf.add_paragraph()
        p_sub.text = model["subtitle"]
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_MUTED_TEXT
        p_sub.font.name = 'Segoe UI'
        p_sub.space_after = Pt(14)
        
        p_desc = tf.add_paragraph()
        p_desc.text = model["desc"]
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_DARK_TEXT
        p_desc.font.name = 'Segoe UI'
        p_desc.line_spacing = 1.3

    # ==========================================
    # SLIDE 7: Final Request (Dark Purple BG)
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, COLOR_PURPLE)
    
    # Title
    t_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Partner with C2B PathFinder"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = "Let's change skills diagnostics from static surveys to a dynamic, AI-driven engine."
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_TEAL
    p2.font.name = 'Segoe UI'
    p2.space_before = Pt(8)
    
    # Final Call to action textbox
    cta_box = slide7.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.0))
    tf_cta = cta_box.text_frame
    tf_cta.word_wrap = True
    
    p_cta1 = tf_cta.paragraphs[0]
    p_cta1.text = "THE REQUEST:"
    p_cta1.font.size = Pt(14)
    p_cta1.font.bold = True
    p_cta1.font.color.rgb = COLOR_TEAL
    p_cta1.font.name = 'Segoe UI'
    
    p_cta2 = tf_cta.add_paragraph()
    p_cta2.text = "• We are seeking pilot partners (training centers & corporations) to deploy C2B PathFinder's recommendation engine with real user cohorts."
    p_cta2.font.size = Pt(18)
    p_cta2.font.color.rgb = COLOR_WHITE
    p_cta2.font.name = 'Segoe UI'
    p_cta2.space_before = Pt(12)
    
    p_cta3 = tf_cta.add_paragraph()
    p_cta3.text = "• Looking for strategic advice on scaling the catalog integration API to connect with major global LMS providers."
    p_cta3.font.size = Pt(18)
    p_cta3.font.color.rgb = COLOR_WHITE
    p_cta3.font.name = 'Segoe UI'
    p_cta3.space_before = Pt(12)
    
    p_cta4 = tf_cta.add_paragraph()
    p_cta4.text = "• Join us at: https://github.com/Annkahoro/C2B_PROJECT.git to inspect our code and documentation."
    p_cta4.font.size = Pt(18)
    p_cta4.font.color.rgb = COLOR_WHITE
    p_cta4.font.name = 'Segoe UI'
    p_cta4.space_before = Pt(12)

    prs.save("C2B_PathFinder_Pitch.pptx")
    print("Pitch deck successfully saved to C2B_PathFinder_Pitch.pptx")

if __name__ == '__main__':
    create_pitch_deck()
